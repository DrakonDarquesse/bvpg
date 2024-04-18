from enum import Enum
import string
from pptx import Presentation
from pptx.enum.dml import MSO_FILL
from bible_passage_api import CuvsPassageApi, KjvPassageApi
from main import BibleBook, Passage, Verse
from jinja2 import Environment, meta

# Configure Jinja2 environment with custom delimiters
env = Environment(
    block_start_string='<=',
    block_end_string='=>',
    variable_start_string='@',
    variable_end_string='@'
)


class BibleBookVerbose(str, Enum):
    genesis = "创世记"
    exodus = "出埃及记"
    proverbs = "箴言"
    romans = "罗马书"


def format_overline(passage: Passage):
    '''
    Format the overline to be in the format of [book] [chapter]:[start_verse]-[end_verse] or etc
    '''
    book = BibleBookVerbose[passage.book.name].value
    if passage.start_verse.chapter == passage.end_verse.chapter:
        return f'{book} {passage.start_verse.chapter}:{passage.start_verse.verse}-{passage.end_verse.verse}'
    else:
        return f'{book} {passage.start_verse.chapter}:{passage.start_verse.verse}-{passage.end_verse.chapter}:{passage.end_verse.verse}'


class PresentationBuilder:
    template = None
    base_file = None
    save_file_name: str | None = None

    def save_file(self):
        raise NotImplementedError("Should have implemented this")

    def build_slide(self):
        raise NotImplementedError("Should have implemented this")

    def get_context(self):
        raise NotImplementedError("Should have implemented this")


class BibleVersePresentationBuilder(PresentationBuilder):
    template = Presentation("test.pptx")

    def __init__(self, save_file_name: str = 'slide.pptx', base_name: str | None = None) -> None:
        self.save_file_name = save_file_name
        self.base_file = Presentation(base_name)

    def get_passages(self):
        passage = Passage(book=BibleBook.romans, start_verse=Verse(
            chapter=15, verse=22), end_verse=Verse(chapter=15, verse=33))
        return [passage]

    def get_slide_names(self):
        slide_names: list[str] = ["verse"]
        passages = self.get_passages()
        if len(passages) > 1:
            slide_names.append('cover_two')
        else:
            slide_names.append('cover_one')
        return slide_names

    # TODO: can let user choose to use slide title or slide internal name, each their own function maybe
    def get_slides(self):
        print("getting slides")
        slide_names = self.get_slide_names()
        return list(filter(lambda x: self.get_title(x) in slide_names, self.template.slides))

    def cover_one(self):
        passage = self.get_passages()[0]
        # TODO: if else check start and end chapter same or not, park it under pasasge model
        return [{
            'cover_one': '读经',
            'chi_book': BibleBookVerbose[passage.book.name].value,
            'eng_book': passage.book.name,
            'verse': f'{passage.start_verse.chapter}: {passage.start_verse.verse}-{passage.end_verse.verse}'
        }]

    def verse(self):
        slide_contexts = []
        passages = self.get_passages()
        for passage in passages:
            english_passage_api = KjvPassageApi()
            chinese_passage_api = CuvsPassageApi()

            english_verses = english_passage_api.retrieve_passage(
                passage)
            chinese_verses = chinese_passage_api.retrieve_passage(
                passage)
            overline = format_overline(passage=passage)

            for i, verse in enumerate(english_verses):
                slide_contexts.append({
                    'verse': overline,
                    'verse_num': verse[0],
                    'chi_verse': chinese_verses[i][1],
                    'eng_verse': verse[1]
                })
        return slide_contexts

    def get_placeholders(self, s: str):
        return list(meta.find_undeclared_variables(env.parse(s)))

    def get_title(self, slide) -> str | None:
        if not slide.shapes.title:
            return None
        placeholders = self.get_placeholders(slide.shapes.title.text)
        if len(placeholders) == 1:
            return placeholders[0]
        elif len(placeholders) < 1:
            # TODO: either cannot multiple words or split and join
            return slide.shapes.title.text.strip().replace(" ", "_")
        else:
            # TODO: either cannot multiple words or split and join
            return "_".join(placeholders)

    def build(self):
        """
        get slide templates and use them to add new slide to base file
        """

        print("start building slides")
        slides = self.get_slides()
        for slide in slides:
            title = self.get_title(slide)
            if title:
                context = getattr(self, title)()
                self.build_slide(slide, context)

    def build_slide(self, template_slide, slide_contexts: list[dict] = []):
        """
        copy background (if its solid color) and all texts along with stylings.
        """
        for slide_config in slide_contexts:
            # TODO: uncouple the slide_layouts

            new_slide = self.base_file.slides.add_slide(
                self.base_file.slide_layouts[6])

            if not template_slide.follow_master_background:
                if template_slide.background.fill.type == MSO_FILL.SOLID:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = template_slide.background.fill.fore_color.rgb
            for shape in template_slide.shapes:
                if not shape.has_text_frame:
                    continue

                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width,
                                                         shape.height)  # add a new textbox to the slide
                new_shape.text_frame.word_wrap = True
                new_shape.text_frame.vertical_anchor = shape.text_frame.vertical_anchor
                for paragraph in shape.text_frame.paragraphs:
                    new_paragraph = new_shape.text_frame.paragraphs[len(
                        new_shape.text_frame.paragraphs)-1]
                    new_paragraph.alignment = paragraph.alignment

                    for run in paragraph.runs:
                        # key = run.text.strip()
                        # text = slide_config[key] if key in slide_config else ''
                        new_run = new_paragraph.add_run()
                        new_run.text = env.from_string(
                            run.text).render(slide_config)
                        new_run.font.name = run.font.name
                        new_run.font.size = run.font.size
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        new_run.font.underline = run.font.underline
                        new_run.font._element.set(
                            'baseline', run.font._element.get('baseline', '0'))
                        if run.font.color.rgb is not None:
                            new_run.font.color.rgb = run.font.color.rgb

                    if len(new_shape.text_frame.paragraphs) != len(shape.text_frame.paragraphs):
                        new_shape.text_frame.add_paragraph()

    def save_file(self):
        self.base_file.save(self.save_file_name)


class ResponsivePresentationBuilder(PresentationBuilder):
    template = Presentation("responsive_reading.pptx")

    def __init__(self, save_file_name: str = 'slide.pptx', base_name: str | None = None) -> None:
        self.save_file_name = save_file_name
        self.base_file = Presentation(base_name)

    def get_passages(self):
        passage = Passage(book=BibleBook.proverbs, start_verse=Verse(
            chapter=13, verse=1), end_verse=Verse(chapter=13, verse=12))
        return [passage]

    def get_slide_names(self):
        slide_names: list[str] = ["cover", "responsive"]
        passage = self.get_passages()[0]
        #  TODO: if pasasge is odd, append "responsive_end"
        if (passage.end_verse.verse - passage.start_verse.verse + 1) % 2 == 1:
            slide_names.append('responsive_end')
        return slide_names

    # TODO: can let user choose to use slide title or slide internal name, each their own function maybe
    def get_slides(self):
        print("getting slides")
        slide_names = self.get_slide_names()
        return list(filter(lambda x: self.get_title(x) in slide_names, self.template.slides))

    def get_placeholders(self, s: str):
        return list(meta.find_undeclared_variables(env.parse(s)))

    def get_title(self, slide) -> str | None:
        if not slide.shapes.title:
            return None
        placeholders = self.get_placeholders(slide.shapes.title.text)
        if len(placeholders) == 1:
            return placeholders[0]
        elif len(placeholders) < 1:
            # TODO: either cannot multiple words or split and join
            return slide.shapes.title.text.strip().replace(" ", "_")
        else:
            # TODO: either cannot multiple words or split and join
            return "_".join(placeholders)

    def build(self):
        """
        get slide templates and use them to add new slide to base file
        """

        print("start building slides")
        slides = self.get_slides()
        for slide in slides:
            title = self.get_title(slide)
            if title:
                context = getattr(self, title)()
                self.build_slide(slide, context)

    def build_slide(self, template_slide, slide_contexts: list[dict] = []):
        """
        copy background (if its solid color) and all texts along with stylings.
        """
        for slide_config in slide_contexts:
            # TODO: uncouple the slide_layouts

            new_slide = self.base_file.slides.add_slide(
                self.base_file.slide_layouts[6])

            if not template_slide.follow_master_background:
                if template_slide.background.fill.type == MSO_FILL.SOLID:
                    new_slide.background.fill.solid()
                    new_slide.background.fill.fore_color.rgb = template_slide.background.fill.fore_color.rgb
            for shape in template_slide.shapes:
                if not shape.has_text_frame:
                    continue

                new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width,
                                                         shape.height)  # add a new textbox to the slide
                new_shape.text_frame.word_wrap = True
                new_shape.text_frame.vertical_anchor = shape.text_frame.vertical_anchor
                for paragraph in shape.text_frame.paragraphs:
                    new_paragraph = new_shape.text_frame.paragraphs[len(
                        new_shape.text_frame.paragraphs)-1]
                    new_paragraph.alignment = paragraph.alignment
                    for run in paragraph.runs:
                        print(run.text)

                    for run in paragraph.runs:
                        print(run.text)
                        new_run = new_paragraph.add_run()
                        new_run.text = env.from_string(
                            run.text).render(slide_config)
                        new_run.font.name = run.font.name
                        new_run.font.size = run.font.size
                        new_run.font.bold = run.font.bold
                        new_run.font.italic = run.font.italic
                        new_run.font.underline = run.font.underline
                        new_run.font._element.set(
                            'baseline', run.font._element.get('baseline', '0'))
                        if run.font.color.rgb is not None:
                            new_run.font.color.rgb = run.font.color.rgb

                    if len(new_shape.text_frame.paragraphs) != len(shape.text_frame.paragraphs):
                        new_shape.text_frame.add_paragraph()

    def save_file(self):
        self.base_file.save(self.save_file_name)

    def cover(self):
        passage = self.get_passages()[0]
        return [{
            'cover': '启应经文 Responsive Reading',
            'chi_book': BibleBookVerbose[passage.book.name].value,
            'eng_book': passage.book.name,
            'chapter_verse': f'{passage.start_verse.chapter}: {passage.start_verse.verse}-{passage.end_verse.verse}'
        }]

    def responsive(self):
        slide_contexts = []
        passage = self.get_passages()[0]
        english_passage_api = KjvPassageApi()
        chinese_passage_api = CuvsPassageApi()

        english_verses = english_passage_api.retrieve_passage(
            passage)
        chinese_verses = chinese_passage_api.retrieve_passage(
            passage)

        for i in range(0, len(english_verses), 2):
            slide_contexts.append({
                'chi_open': chinese_verses[i][1],
                'eng_open': english_verses[i][1],
                'chi_close': chinese_verses[i+1][1],
                'eng_close': english_verses[i+1][1],
                'responsive': ''
            })
        return slide_contexts

    def responsive_end(self):
        passage = self.get_passages()[0]
        english_passage_api = KjvPassageApi()
        chinese_passage_api = CuvsPassageApi()

        english_verses = english_passage_api.retrieve_passage(
            passage)
        chinese_verses = chinese_passage_api.retrieve_passage(
            passage)
        last_index = len(english_verses)

        return {
            'chi_close': chinese_verses[last_index][1],
            'eng_close': english_verses[last_index][1],
            'responsive_end': ''
        }


slide = ResponsivePresentationBuilder(base_name="base_wide.pptx")
slide.build()
slide.save_file()
