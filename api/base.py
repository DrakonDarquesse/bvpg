from enum import Enum
from pptx import Presentation as ppt
from pptx.presentation import Presentation, Slides
from pptx.enum.dml import MSO_FILL  # type: ignore
from environment import env


def duplicate_shape(source_shape, target_shape):
    """
    Args:
        source_shape  
        target_shape

    copies properties of source shape into target shape, including paragraphs and text styles
    """
    target_shape.left = source_shape.left
    target_shape.top = source_shape.top
    target_shape.width = source_shape.width
    target_shape.height = source_shape.height

    target_shape.text_frame.word_wrap = True
    target_shape.text_frame.vertical_anchor = source_shape.text_frame.vertical_anchor

    for paragraph in source_shape.text_frame.paragraphs:
        # when iterating paragraph, update the last paragraph in new_shape.text_frame
        new_paragraph = target_shape.text_frame.paragraphs[len(
            target_shape.text_frame.paragraphs)-1]
        new_paragraph.alignment = paragraph.alignment

        for run in paragraph.runs:
            new_run = new_paragraph.add_run()
            new_run.text = run.text
            new_run.font.name = run.font.name
            new_run.font.size = run.font.size
            new_run.font.bold = run.font.bold
            new_run.font.italic = run.font.italic
            new_run.font.underline = run.font.underline
            new_run.font._element.set(
                'baseline', run.font._element.get('baseline', '0'))
            if run.font.color.rgb is not None:
                new_run.font.color.rgb = run.font.color.rgb

        if len(target_shape.text_frame.paragraphs) != len(source_shape.text_frame.paragraphs):
            target_shape.text_frame.add_paragraph()


def duplicate_slide(source_slide, target_slide):
    """
    copy background (if its solid color) and all texts along with stylings.
    """

    # set background is its solid fill
    try:

        if not source_slide.follow_master_background:
            if source_slide.background.fill.type == MSO_FILL.SOLID:
                target_slide.background.fill.solid()
                target_slide.background.fill.fore_color.rgb = source_slide.background.fill.fore_color.rgb
    except:
        pass

    title = source_slide.shapes.title

    duplicate_shape(title, target_slide.shapes.title)
    for shape in source_slide.shapes:
        if not shape.has_text_frame:
            continue

        if shape == title:
            continue

        new_shape = target_slide.shapes.add_textbox(shape.left, shape.top, shape.width,
                                                    shape.height)  # add a new textbox to the slide

        duplicate_shape(shape, new_shape)


def render_slide_data(slide, context, env):

    slide.shapes.title.text_frame.paragraphs[0].runs[0].text = context['data'].get(
        'title', slide.shapes.title.text)
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:

            for run in paragraph.runs:
                run.text = env.from_string(
                    run.text).render(context['data'])


class PresentationTemplate:

    def __init__(self, template: Presentation) -> None:
        if not isinstance(template, Presentation):
            raise Exception("Wrong template")
        self.template = template
        self.slides = self.get_titled_slides()

    def get_titled_slides(self):
        """
        only returns slide with title
        """
        print("getting slides")

        template_slides: Slides = self.template.slides  # type: ignore
        return list(filter(lambda slide: bool(slide.shapes.title), template_slides))

    def get_slide_by_title(self, title):
        try:
            return next(
                (slide for slide in self.slides if slide.shapes.title.text == title))
        except:
            raise Exception(f"no slide with the title {title}")


class BlankPresentation:
    class SlideLayout(int, Enum):
        TITLE_SLIDE = 5

    def __init__(self, save_file_name: str | None = None) -> None:
        self.base_file = ppt("template/base_wide.pptx")
        self.save_file_name = save_file_name

    def save_file(self, file_name=None):
        self.base_file.save(file_name or self.save_file_name)

    def new_title_slide(self):
        base_file = self.base_file
        return base_file.slides.add_slide(
            base_file.slide_layouts[self.SlideLayout.TITLE_SLIDE])

    def get_base_file(self):
        return self.base_file


class Context:

    def get_context(self):
        raise NotImplementedError("Should have implemented this")


class PresentationBuilder():

    def __init__(self, presentation: BlankPresentation, template: PresentationTemplate, context: Context) -> None:
        self.presentation = presentation
        self.template = template
        self.context = context

    def build(self):
        """
        get slide templates and use them to add new slide to base file
        """

        print("start building slides")
        contexts = self.context.get_context()
        for context in contexts:
            self.build_slide(context)

    def build_slide(self, context):
        template_slide = self.template.get_slide_by_title(context["title"])
        new_slide = self.presentation.new_title_slide()

        duplicate_slide(template_slide, new_slide)
        render_slide_data(new_slide, context, env)
