from parse_esv_json import ESVJSONParser
from fastapi import FastAPI, Response
from pptx import Presentation
from io import BytesIO

from bible_passage_api import CuvsPassageApi, EsvPassageApi, KjvPassageApi
from data.bible import BibleBook
from base import BlankPresentation, PresentationBuilder, PresentationTemplate
from responsive_context import ResponsiveContext
from passage_context import PassageContext
from models import Passage, Verse
from mongodb import mongodb_session

app = FastAPI()


@app.post("/bvpg/slides/bible-reading", response_class=Response)
def build_slide_bible_reading(passages: list[Passage]):
    output = BytesIO()
    slide = PresentationBuilder(
        presentation=BlankPresentation(), template=PresentationTemplate(
            Presentation("template/passage_template.pptx")), context=PassageContext(passages))

    slide.build()
    slide.presentation.save_file(output)
    return Response(output.getvalue(), media_type='application/pptx')


@app.post("/bvpg/slides/responsive-reading", response_class=Response)
async def build_slide_responsive_reading(passages: list[Passage]):
    output = BytesIO()
    slide = PresentationBuilder(
        presentation=BlankPresentation(), template=PresentationTemplate(
            Presentation("template/responsive_reading.pptx")), context=ResponsiveContext(passages))
    slide.build()
    slide.presentation.save_file(output)
    return Response(output.getvalue(), media_type='application/pptx')


@app.post("/bvpg/slides/passages", response_model=list[Passage])
async def passages(passages: list[Passage]) -> list[Passage]:
    return passages


@app.get("/bvpg/api/passage/rcuvss/books/{book}/chapters/{chapter}/verses/{start_verse}/{end_verse}")
async def rcuvss_passage(book: BibleBook, chapter: int, start_verse: int, end_verse: int):
    passage = Passage(
        book=book,
        start_verse=Verse(chapter=chapter, verse=start_verse),
        end_verse=Verse(chapter=chapter, verse=end_verse)
    )

    text = CuvsPassageApi().retrieve_passage(
        passage=passage
    )

    return text

@app.get("/bvpg/api/passage/esv/books/{book}/chapters/{chapter}/verses/{start_verse}/{end_verse}")
async def esv_passage(book: BibleBook, chapter: int, start_verse: int, end_verse: int):
    passage = Passage(
        book=book,
        start_verse=Verse(chapter=chapter, verse=start_verse),
        end_verse=Verse(chapter=chapter, verse=end_verse)
    )


    text = EsvPassageApi().retrieve_passage(
        passage=passage
    )
    print(text)
    return text



@app.get("/bvpg/api/passage/kjv/books/{book}/chapters/{chapter}/verses/{start_verse}/{end_verse}")
async def kjv_passage(book: BibleBook, chapter: int, start_verse: int, end_verse: int):
    passage = Passage(
        book=book,
        start_verse=Verse(chapter=chapter, verse=start_verse),
        end_verse=Verse(chapter=chapter, verse=end_verse)
    )

    text = KjvPassageApi().retrieve_passage(
        passage=passage
    )

    return text


@app.get("/")
async def root():
    return {"message": "Hello World"}


# this function is used to upload bible verses into mongodb
# not needed in operation
async def upload():
    esv_json_parser = ESVJSONParser("../../mdbible/json/ESV.json", mongodb_session)
    esv_json_parser.load_data()
    esv_json_parser.extract_scriptures()
    esv_json_parser.upload_bible_to_database("esv", "verses")
