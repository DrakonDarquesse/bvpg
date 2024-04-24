import re
import time
from pptx import Presentation
# from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_FILL
import requests
import sqlite3


def chinese_bible_api(book, start_chapter, end_chapter, start_verse, end_verse):

    try:

        # Connect to DB and create a cursor
        sqliteConnection = sqlite3.connect('bible_cuvs.db')
        cursor = sqliteConnection.cursor()
        print('DB Init')

        # Write a query and execute it with cursor
        query = \
            f'''
            WITH book AS 
            (SELECT bible.ID, Bible.ChapterSN, Bible.VerseSN, Bible.Lection FROM Bible
            JOIN BibleID ON Bible.VolumeSN = BibleID.SN 
            WHERE Bible.VolumeSN = {book})

            SELECT Bible.VerseSN, Bible.Lection FROM Bible
            WHERE BIBLE.ID BETWEEN
            (
                SELECT ID FROM book
                WHERE ChapterSN = {start_chapter} AND VerseSN = {start_verse}

            ) AND
            (
                SELECT ID FROM book
                WHERE ChapterSN = {end_chapter} AND VerseSN = {end_verse}

            );
            '''
        cursor.execute(query)

        # Fetch and output result
        result = cursor.fetchall()

        # Close the cursor
        cursor.close()

        return result

    # Handle errors
    except sqlite3.Error as error:
        print('Error occurred - ', error)

    # Close DB Connection irrespective of success
    # or failure
    finally:

        if sqliteConnection:
            sqliteConnection.close()
            print('SQLite Connection closed')

    return []


def match_number_and_verse_pair(passage):
    '''
    Match the entire pattern of bracket number follow by text
    '''
    verses = re.findall(r'\[(\d+)\](.*?)(?=\[|$)', passage, re.DOTALL)
    # Filter out empty strings and whitespace
    verses = [(verse[0], verse[1].strip())
              for verse in verses if verse[1].strip()]
    return verses


def split_verse_and_match_number(passage):
    '''
    Split the passage using verse numbers as delimiter
    and match the verse number with the verse text
    '''

    # Find all verse numbers in the passage
    pattern = r'(?<=\[)\d+(?=\])'
    verse_numbers = [number for number in re.findall(pattern, passage)]

    # Split the text by bracket numbers
    verses = re.split(r'\[\d+\]', passage)

    # Filter out empty strings and whitespace
    verses = [string.strip() for string in verses if string.strip()]

    if len(verses) != len(verse_numbers):
        print('Verse numbers and verses do not match')

    return list(zip(verse_numbers, verses))


def english_bible_api(book, start_chapter, end_chapter, start_verse, end_verse):
    bible = 'de4e12af7f28f599-01'

    passage = f'{book}.{start_chapter}.{start_verse}-{book}.{end_chapter}.{end_verse}'

    url = f'https://api.scripture.api.bible/v1/bibles/{bible}/passages/{passage}'

    r = requests.get(url, headers={
        'api-key': 'f89e0fc2938f4d054717716279057d45', 'accept': 'application/json'}, params={'content-type': 'text', 'include-titles': 'false', 'include-verse-numbers': 'true'})

    content = r.json().get('data').get('content')

    return match_number_and_verse_pair(content)


def format_overline(book, chapter, start_verse, end_verse):
    '''
    Format the overline to be in the format of [book] [chapter]:[start_verse]-[end_verse]
    '''
    return f'{book} {chapter}:{start_verse}-{end_verse}'


def create_slide_config():
    english_verses = english_bible_api('EXO', 9, 10, 34, 5)
    chinese_verses = chinese_bible_api(2, 9, 10, 34, 5)
    overline = '马太福音 9:9-17'

    slide_configs = []
    for i, verse in enumerate(english_verses):
        slide_configs.append({
            'overline': overline,
            'verse': verse[0],
            'chi_text': chinese_verses[i][1],
            'eng_text': verse[1]
        })
    return slide_configs


slide_configs = create_slide_config()

'''
create a abstraction layer for getting the bible verses

for each api, create a class or something that inherits from the abstraction layer

custom template, let user load the file, then load the template context, then let user fill in the form
'''

# slide_config = {
#     'overline': ['读经'],
#     'title': ['创世记', 'Genesis', '1: 12']
# }

# slide_config = {
#     'overline': [['马太福音 9:9-17']],
#     'passage': [
#         ['9', '耶稣从那里往前走，看见一个人，名叫马太，坐在税关上，就对他说：“你跟从我来！”他就起来跟从了耶稣。'],
#         [''],
#         ['9', 'And as Jesus passed forth from thence, he saw a man, named Matthew, sitting at the receipt of custom: and he saith unto him, Follow me. And he arose, and followed him.']
#     ]
# }

# slide_config = {
#     'overline': '马太福音 9:9-17',
#     'verse': '9',
#     'chi_text': '耶稣从那里往前走，看见一个人，名叫马太，坐在税关上，就对他说：“你跟从我来！”他就起来跟从了耶稣。',
#     'eng_text': 'And as Jesus passed forth from thence, he saw a man, named Matthew, sitting at the receipt of custom: and he saith unto him, Follow me. And he arose, and followed him.'
# }


template = Presentation('test.pptx')
prs = Presentation('test1.pptx')
slide = template.slides[1]

for slide_config in slide_configs:
    # loop thru the shapes and match runs text with slide_config
    new_slide = prs.slides.add_slide(prs.slide_layouts[6])

    if not slide.follow_master_background:
        if slide.background.fill.type == MSO_FILL.SOLID:
            new_slide.background.fill.solid()
            new_slide.background.fill.fore_color.rgb = slide.background.fill.fore_color.rgb
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width,
                                                 shape.height)  # add a new textbox to the slide
        new_shape.text_frame.word_wrap = True
        for paragraph in shape.text_frame.paragraphs:
            new_paragraph = new_shape.text_frame.paragraphs[len(
                new_shape.text_frame.paragraphs)-1]
            new_paragraph.alignment = paragraph.alignment

            for run in paragraph.runs:
                key = run.text.strip()
                print(len(key))
                text = slide_config[key] if key in slide_config else ''
                new_run = new_paragraph.add_run()
                new_run.text = text
                new_run.font.name = run.font.name
                new_run.font.size = run.font.size
                new_run.font.bold = run.font.bold
                new_run.font.italic = run.font.italic
                new_run.font.underline = run.font.underline
                new_run.font._element.set(
                    'baseline', run.font._element.get('baseline', '0'))
                if run.font.color.rgb is not None:
                    new_run.font.color.rgb = run.font.color.rgb

            if len(new_shape.text_frame.paragraphs) != len(shape.text_frame.paragraphs):
                new_shape.text_frame.add_paragraph()

# # loop thru the slide_config dictionary of shape names
# for slide_config in slide_configs:
#     for key, value in slide_config.items():
#         shape = next(
#             (shape for shape in slide.shapes if shape.name == key), None)

#         if shape is None:
#             continue

#         if not shape.has_text_frame:
#             continue

#         new_shape = new_slide.shapes.add_textbox(
#             shape.left, shape.top, shape.width, shape.height)  # add a new textbox to the slide
#         new_shape.text_frame.word_wrap = True

#         for i, paragraph in enumerate(shape.text_frame.paragraphs):
#             new_paragraph = new_shape.text_frame.paragraphs[len(
#                 new_shape.text_frame.paragraphs)-1]
#             new_paragraph.alignment = paragraph.alignment

#             for j, run in enumerate(paragraph.runs):
#                 new_run = new_paragraph.add_run()
#                 print(run.text)
#                 print(value[i][j])
#                 new_run.text = value[i][j]
#                 new_run.font.name = run.font.name
#                 new_run.font.size = run.font.size
#                 new_run.font.bold = run.font.bold
#                 new_run.font.italic = run.font.italic
#                 new_run.font.underline = run.font.underline
#                 new_run.font._element.set(
#                     'baseline', run.font._element.get('baseline', '0'))
#                 if run.font.color.rgb is not None:
#                     new_run.font.color.rgb = run.font.color.rgb

#             if len(new_shape.text_frame.paragraphs) != len(shape.text_frame.paragraphs):
#                 new_shape.text_frame.add_paragraph()


# # Copy the  of the original slide to the new
# # loop thru the shapes in the tmeplate slides
# for shape in slide.shapes:
#     if not shape.has_text_frame:
#         continue

#     new_shape = new_slide.shapes.add_textbox(shape.left, shape.top, shape.width,
#                                              shape.height)  # add a new textbox to the slide
#     new_shape.text_frame.word_wrap = True
#     for paragraph in shape.text_frame.paragraphs:
#         new_paragraph = new_shape.text_frame.paragraphs[len(
#             new_shape.text_frame.paragraphs)-1]
#         new_paragraph.alignment = paragraph.alignment

#         for run in paragraph.runs:
#             print(run.text)
#             new_run = new_paragraph.add_run()
#             new_run.text = run.text
#             new_run.font.name = run.font.name
#             new_run.font.size = run.font.size
#             new_run.font.bold = run.font.bold
#             new_run.font.italic = run.font.italic
#             new_run.font.underline = run.font.underline
#             if run.font.color.rgb is not None:
#                 new_run.font.color.rgb = run.font.color.rgb

#         if len(new_shape.text_frame.paragraphs) != len(shape.text_frame.paragraphs):
#             new_shape.text_frame.add_paragraph()

prs.save('test1.pptx')
